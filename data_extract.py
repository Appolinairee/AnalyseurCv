import os

from pdfminer.high_level import extract_text
from pptx import Presentation
from docx import Document


def extract_text_from_txt(file_path):
    with open(file_path, 'r') as f:
        text = f.read()
    return text

def extract_text_from_pdf(file_path):
    try:
        text = extract_text(file_path)
        return text
    except Exception as e:
        print(f"Une erreur s'est produite \n{e}")


def extract_text_from_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += f"{shape.text}\n"
        return text
    except Exception as e:
        print(f"Une erreur s'est produite \n{e}")


def extract_text_from_docx(file_path):
    text = ""
    try:
        doc = Document(file_path)
        for paragraph in doc.paragraphs:
            text += f"{paragraph.text}\n"
        return text
    except Exception as e:
        print(f"Une erreur s'est produite \n{e}")

def extract_text_from_img(file_path):
    # Not implemented
    return ""


def extract_text_from_file(file_path):
    # récupérer l'extension du fichier
    ext = file_path.split('.')[-1] # le dernier élément est l'extension
    ppt = ['pptx', 'ppt']
    word = ['docx', 'doc']
    img = ['jpg', 'jpeg', 'png', 'bmp', 'webp']
    pdf = ['pdf']
    txt = ['txt']
    
    if ext not in ppt + word + img + pdf:
        print("Type non accepté")
        return -1
    else:
        if ext in ppt:
            return extract_text_from_pptx(file_path)
        elif ext in word:
            return extract_text_from_docx(file_path)
        elif ext in pdf:
            return extract_text_from_pdf(file_path)
        elif ext in txt:
            return extract_text_from_txt(file_path)
        else:
            return extract_text_from_img(file_path)


def extract_text_from_many(lst: list, dico=False):
    # retourner une liste ou un dictionnaire
    lst_texts = [extract_text_from_file(file_path) for file_path in lst]
    dico_texts = {file_path: extract_text_from_file(file_path) 
                  for file_path in lst}
    return (dico_texts if dico else lst_texts)


if __name__ == '__main__':
    os.chdir("./assets/CV_Appolinaire_ADANDE (1).pdf")
    lst_file_paths = list(el for el in os.listdir())
    print(f"ENTREE \n{lst_file_paths}\n")
    # extraire le texte de chaque fichier
    lst_extracted_texts = extract_text_from_many(lst_file_paths)
    # on peut utiliser dico = True pour obtenir un dictionnaire en sortie
    print(f"SORTIE \n{lst_extracted_texts}")  


# Remarques
# La lecture des tableaux n'est pas faite avec les documents words
# La présence d'images n'est un problème pour aucune des bibliothèques utlisées
# 
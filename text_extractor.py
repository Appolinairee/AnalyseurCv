import os
import fitz  # pymupdf
from pptx import Presentation
from docx import Document
import pytesseract
from PIL import Image

pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'


class TextExtractor:
    def __init__(self, file_path):
        self.file_path = file_path

    def extract_from_pdf(self):
        try:
            with fitz.open(self.file_path) as pdf_document:
                texte_complet = ""
                for page_num in range(pdf_document.page_count):
                    page = pdf_document[page_num]
                    texte_page = page.get_text("text")
                    texte_complet += texte_page
                return texte_complet
        except Exception as e:
            print(f"Erreur lors de l'extraction du texte du PDF : {e}")
            return None

    def extract_from_doc(self):
        try:
            doc = Document(self.file_path)
            texte_complet = ""
            for paragraphe in doc.paragraphs:
                texte_complet += paragraphe.text + "\n"
            return texte_complet.strip()
        except Exception as e:
            print(f"Erreur lors de l'extraction du texte du document : {e}")
            return None

    def extract_from_ppt(self):
        try:
            presentation = Presentation(self.file_path)
            texte_complet = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texte_complet += shape.text + "\n"
            return texte_complet.strip()
        except Exception as e:
            print(f"Erreur lors de l'extraction du texte de la présentation : {e}")
            return None

    def extract_from_txt(self):
        try:
            with open(self.file_path, 'r', encoding='utf-8') as fichier:
                contenu = fichier.read()
                return contenu
        except Exception as e:
            print(f"Erreur lors de la lecture du fichier texte : {e}")
            return None

    def extract_from_image(self):
        try:
            image = Image.open(self.file_path)
            texte_extrait = pytesseract.image_to_string(image, lang='eng')
            return texte_extrait
        except Exception as e:
            print(f"Erreur lors de l'extraction du texte de l'image : {e}")
            return None

    def extract(self):
        # determination of file extension
        extension = os.path.splitext(self.file_path)[1][1:]
        extension = extension.lower()

        if extension == "pdf":
            return self.extract_from_pdf()

        elif extension == "doc" or extension == "docx":
            return self.extract_from_doc()

        elif extension == "ppt" or extension == "pptx":
            return self.extract_from_ppt()

        elif extension == "txt":
            return self.extract_from_txt()

        elif extension in ["jpg", "jpeg", "png"]:
            return self.extract_from_image()

        else:
            print("Extension non prise en compte")
            return None

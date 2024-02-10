import os
import fitz 
from pptx import Presentation
from docx import Document
import pytesseract
from PIL import Image

pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'


def extract_from_pdf(file_path):
    try:
        with fitz.open(file_path) as pdf_document:
            texte_complet = ""
            
            for page_num in range(pdf_document.page_count):
                page = pdf_document[page_num]
                texte_page = page.get_text("text")
                
                texte_complet += texte_page
            
            return texte_complet

    except Exception as e:
        print(f"Erreur lors de l'extraction du texte : {e}")
        return None



def extract_from_doc(file_path):
    try:

        doc = Document(file_path)

        texte_complet = ""

        for paragraphe in doc.paragraphs:
            texte_complet += paragraphe.text + "\n"

        return texte_complet.strip()

    except Exception as e:
        print(f"Erreur lors de l'extraction du texte : {e}")
        return None


def extract_from_ppt(chemin_pptx):
    try:
        presentation = Presentation(chemin_pptx)

        texte_complet = ""

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texte_complet += shape.text + "\n"

        return texte_complet.strip()

    except Exception as e:
        print(f"Erreur lors de l'extraction du texte : {e}")
        return None



def  extract_from_txt(chemin_fichier):
    try:
        with open(chemin_fichier, 'r', encoding='utf-8') as fichier:
            contenu = fichier.read()

            return contenu

    except Exception as e:
        print(f"Erreur lors de la lecture du fichier : {e}")
        return None
    

def extract_from_image(chemin_image):
    try:
        image = Image.open(chemin_image)

        texte_extrait = pytesseract.image_to_string(image, lang='eng')

        return texte_extrait

    except Exception as e:
        print(f"Erreur lors de l'extraction du texte de l'image : {e}")
        return None


def extract(file_path):
    # determination of file extension
    extension = os.path.splitext(file_path)[1][1:]
    extension = extension.lower()

    if extension == "pdf" :
        return extract_from_pdf(file_path)

    elif extension == "doc" or extension == "docx" :
        return extract_from_doc(file_path)

    elif extension == "ppt" or extension == "pptx" :
        return extract_from_ppt(file_path)

    elif extension == "txt" :
        return extract_from_txt(file_path)

    elif extension == extension in ["jpg", "jpeg", "png"] :
        return extract_from_image(file_path)

    else :
        print("Extension non prise en compte")
    return None


if __name__ == "__main__":
    
    # cv files to extract
    file_paths = ["./assets/calendar.txt", "./assets/CV_Appolinaire_ADANDE (1).pdf", "./assets/Capture d'écran 2024-02-10 125444.png", "./assets/modèle rapports.docx"]

    for file_path in file_paths:
        text = extract(file_path) 
        print(f"Fichier {file_path} \n\n {text} \n\n\n")